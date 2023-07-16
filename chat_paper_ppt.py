import io
import json
import math
import os
import re
import sys
from shutil import copyfile, move

import zhon
import cv2
import pptx
from PIL import Image
from arxiv import arxiv
import fitz
from pathvalidate import sanitize_filepath
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from chat_paper import chat_paper_main, Paper, Reader
import argparse
from sanitize_filename import sanitize
import numpy as np
def analyze_text(text):
    print("Analyze Text:\n" + text)
    # 定义正则表达式模式
    pattern_title = r"[标题|Title]:\s*(.*)"
    pattern_author = r"[作者|Authors]:\s*(.*)"
    pattern_affiliation = r"[机构|Affiliation]:\s*(.*)"
    pattern_keywords = r"[关键词|Keywords]:\s*(.*)"
    pattern_abstract = r"[摘要|Summary]:\s*([\s\S]*)"
    pattern_conclu = r"Conclusion:\s*([\s\S]*)"

    # 使用正则表达式匹配信息
    match_title = re.search(pattern_title, text)
    match_author = re.search(pattern_author, text)
    match_affiliation = re.search(pattern_affiliation, text)
    match_keywords = re.search(pattern_keywords, text)
    match_abstract = re.search(pattern_abstract, text)
    match_conclu = re.search(pattern_conclu, text)

    # 将提取出的信息放入字典中
    result_dict = {
        "标题": match_title.group(1).strip(),
        "作者": match_author.group(1).strip(),
        "单位": match_affiliation.group(1).strip(),
        "关键字": match_keywords.group(1).strip(),
        "摘要": match_abstract.group(1).strip(),
        "总结": match_conclu.group(1).strip()
    }
    # print(result_dict)
    return result_dict,match_title.group(1).strip()

def extract_imgs_from_pdf(title, file):
    # STEP 2
    # file path you want to extract images from

    # open the file
    pdf_file = fitz.open(file)
    # Minimum width and height for extracted images
    min_width = 150
    min_height = 150
    # Output directory for the extracted images

    output_dir = "export/"+sanitize(title)
    if not os.path.exists(output_dir):
        os.mkdir(output_dir)
        os.mkdir(output_dir+"/small")
    # Desired output image format
    output_format = "png"
    # STEP 3
    # iterate over PDF pages
    first_image_page = -1
    first_image_number = -1
    images = []
    for page_index in range(len(pdf_file)):

        # get the page itself
        page = pdf_file[page_index]
        image_list = page.get_images()

        # printing number of images found in this page
        # if image_list:
        #     print(
        #         f"[+] Found a total of {len(image_list)} images in page {page_index}")
        # else:
        #     print("[!] No images found on page", page_index)
        for image_index, img in enumerate(page.get_images(), start=1):
            # get the XREF of the image
            xref = img[0]
            # print("xref", xref)

            # extract the image bytes
            base_image = pdf_file.extract_image(xref)
            image_bytes = base_image["image"]

            # get the image extension
            image_ext = base_image["ext"]
            # Load it to PIL
            image = Image.open(io.BytesIO(image_bytes))
            # print(image_index, image.width, image.height)
            # Check if the image meets the minimum dimensions and save it

            if image.width >= min_width and image.height >= min_height:
                if first_image_page == -1:
                    first_image_page = page_index
                    first_image_number = image_index
                image.convert('RGB').save(
                    open(os.path.join(output_dir, f"image{page_index + 1}_{image_index}.{output_format}"), "wb"),
                    format=output_format.upper())
                images.append(os.path.join(output_dir, f"image{page_index + 1}_{image_index}.{output_format}"))
            else:
                # print(f"[-] Skipping image {image_index} on page {page_index} due to its small size.")
                image.convert('RGB').save(
                    open(os.path.join(output_dir+"/small", f"image{page_index + 1}_{image_index}.{output_format}"), "wb"),
                    format=output_format.upper())
    return images


def write_pptx(title, result_dict, images, spacing=8):

    ppt_name = "export/" + title + ".pptx"
    print("Write PPT File", ppt_name)
    # 打开模板pptx文件
    template = pptx.Presentation('export/template.pptx')
    # 获取第一张幻灯片，共1张
    slide = template.slides[0]
    # 获取文本框对象
    title_box = slide.shapes[0].text_frame
    abstract_box = slide.shapes[1].text_frame
    author_box = slide.shapes[2].text_frame
    affiliation_box = slide.shapes[3].text_frame
    keywords_box = slide.shapes[4].text_frame

    # 填充文本框
    title_box.paragraphs[0].text = result_dict['标题']
    author_box.paragraphs[0].text = result_dict['作者']
    affiliation_box.paragraphs[0].text = result_dict['单位']
    if isinstance(result_dict['关键字'], list):
        keywords_box.paragraphs[0].text = ', '.join(result_dict['关键字'])
    elif isinstance(result_dict['关键字'], str):
        keywords_box.paragraphs[0].text = result_dict['关键字']
    else:
        keywords_box.paragraphs[0].text = ""
    abstract_box.paragraphs[0].text = result_dict['摘要']

    # 设置文本框的对齐方式和自动调整大小
    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    '''
    title_box.paragraphs[0].font.name = '微软雅黑'
    title_box.paragraphs[0].font.size = Pt(20)
    title_box.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    author_box.paragraphs[0].font.name = '微软雅黑'
    author_box.paragraphs[0].font.size = Pt(16)
    author_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    affiliation_box.paragraphs[0].font.name = '微软雅黑'
    affiliation_box.paragraphs[0].font.size = Pt(16)
    affiliation_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    keywords_box.paragraphs[0].font.name = '微软雅黑'
    keywords_box.paragraphs[0].font.size = Pt(12)
    keywords_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    keywords_box.paragraphs[0].font.italic = True
    abstract_box.paragraphs[0].font.name = '微软雅黑'
    abstract_box.paragraphs[0].font.size = Pt(16)
    abstract_box.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    '''
    slide.shapes[2].top = slide.shapes[0].top + slide.shapes[0].height + Pt(spacing)
    slide.shapes[3].top = slide.shapes[2].top + slide.shapes[2].height + Pt(spacing)
    slide.shapes[4].top = slide.shapes[3].top + slide.shapes[3].height + Pt(spacing)
    slide.shapes[1].top = slide.shapes[4].top + slide.shapes[4].height + Pt(spacing)

    pic = slide.shapes.add_picture(images[0],  # io.BytesIO(img_bytes),
                                   left=template.slide_width * 0.55,
                                   top=slide.shapes[1].top,
                                   width=template.slide_width * 0.4)
    image_slide = template.slides.add_slide(template.slide_layout[6])
    image_count = math.sqrt(len(images))
    for index, image in enumerate(images[1:]):
        pic = image_slide.shapes.add_picture(image,#io.BytesIO(img_bytes),
                                       left=template.slide_width * (1/image_count*(index%image_count)),
                                       top=slide.shapes[1].top* (1/image_count*(int(index/image_count))),
                                       width=template.slide_width * (1/image_count))


    # 保存填充后的pptx文件
    template.save(ppt_name)
    return ppt_name

def write_image_pptx(title, result_dict, images, spacing=8):

    ppt_name = "export/" + title + ".pptx"
    print("Write PPT File", ppt_name)
    # 打开模板pptx文件
    # template = pptx.Presentation('export/template.pptx')
    cur_ppt = pptx.Presentation('export/template.pptx')
    # 获取第一张幻灯片，共1张
    # slide_first = template.slides[0]
    sum_slide = cur_ppt.slides.add_slide(cur_ppt.slide_layouts[6])
    # slide = template.slides.add_slide(sum_slide)
    for i in range(5):
        sum_slide.shapes.add_textbox(0,0,cur_ppt.slide_width/2,cur_ppt.slide_height)

    title_box = sum_slide.shapes[0].text_frame
    abstract_box = sum_slide.shapes[1].text_frame
    author_box = sum_slide.shapes[2].text_frame
    affiliation_box = sum_slide.shapes[3].text_frame
    keywords_box = sum_slide.shapes[4].text_frame

    # 填充文本框
    title_box.paragraphs[0].text = result_dict['标题']
    author_box.paragraphs[0].text = result_dict['作者']
    affiliation_box.paragraphs[0].text = result_dict['单位']
    if isinstance(result_dict['关键字'], list):
        keywords_box.paragraphs[0].text = ', '.join(result_dict['关键字'])
    elif isinstance(result_dict['关键字'], str):
        keywords_box.paragraphs[0].text = result_dict['关键字']
    else:
        keywords_box.paragraphs[0].text = ""
    abstract_box.paragraphs[0].text = result_dict['摘要']

    # 设置文本框的对齐方式和自动调整大小
    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    '''
    title_box.paragraphs[0].font.name = '微软雅黑'
    title_box.paragraphs[0].font.size = Pt(20)
    title_box.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    author_box.paragraphs[0].font.name = '微软雅黑'
    author_box.paragraphs[0].font.size = Pt(16)
    author_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    affiliation_box.paragraphs[0].font.name = '微软雅黑'
    affiliation_box.paragraphs[0].font.size = Pt(16)
    affiliation_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    keywords_box.paragraphs[0].font.name = '微软雅黑'
    keywords_box.paragraphs[0].font.size = Pt(12)
    keywords_box.paragraphs[0].font.color.rgb = RGBColor(96, 96, 96)
    keywords_box.paragraphs[0].font.italic = True
    abstract_box.paragraphs[0].font.name = '微软雅黑'
    abstract_box.paragraphs[0].font.size = Pt(16)
    abstract_box.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    '''
    sum_slide.shapes[2].top = sum_slide.shapes[0].top + sum_slide.shapes[0].height + Pt(spacing)
    sum_slide.shapes[3].top = sum_slide.shapes[2].top + sum_slide.shapes[2].height + Pt(spacing)
    sum_slide.shapes[4].top = sum_slide.shapes[3].top + sum_slide.shapes[3].height + Pt(spacing)
    sum_slide.shapes[1].top = sum_slide.shapes[4].top + sum_slide.shapes[4].height + Pt(spacing)

    pic = sum_slide.shapes.add_picture(images[0],  # io.BytesIO(img_bytes),
                                   left=cur_ppt.slide_width * 0.55,
                                   top=sum_slide.shapes[1].top,
                                   width=cur_ppt.slide_width * 0.4)

    # prs = Presentation()
    image_slide = cur_ppt.slides.add_slide(cur_ppt.slide_layouts[6])
    image_count = math.sqrt(len(images))
    for index, image in enumerate(images[1:]):
        # print("left",(1/image_count*(index%image_count)),"slide_height",1/image_count*(int(index/image_count)))
        pic = image_slide.shapes.add_picture(image,#io.BytesIO(img_bytes),
                                       left=cur_ppt.slide_width * (1/image_count*(index%image_count)),
                                       top=cur_ppt.slide_height *  (1/image_count*(int(index/image_count))),
                                       width=cur_ppt.slide_width * (1/image_count))


    # 保存填充后的pptx文件
    cur_ppt.save(ppt_name)
    return ppt_name

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument("--pdf_path", type=str, default=r'demo.pdf',
                        help="if none, the bot will download from arxiv with query")
    # parser.add_argument("--pdf_path", type=str, default=r'C:\Users\Administrator\Desktop\DHER\RHER_Reset\ChatPaper', help="if none, the bot will download from arxiv with query")
    # parser.add_argument("--pdf_path", type=str, default='', help="if none, the bot will download from arxiv with query")
    parser.add_argument("--query", type=str, default='all: ChatGPT robot',
                        help="the query string, ti: xx, au: xx, all: xx,")
    parser.add_argument("--key_word", type=str, default='reinforcement learning',
                        help="the key word of user research fields")
    parser.add_argument("--filter_keys", type=str, default='ChatGPT robot',
                        help="the filter key words, 摘要中每个单词都得有，才会被筛选为目标论文")
    parser.add_argument("--max_results", type=int, default=1, help="the maximum number of results")
    # arxiv.SortCriterion.Relevance
    parser.add_argument("--sort", type=str, default="Relevance", help="another is LastUpdatedDate")
    parser.add_argument("--save_image", default=False,
                        help="save image? It takes a minute or two to save a picture! But pretty")
    parser.add_argument("--file_format", type=str, default='md', help="导出的文件格式，如果存图片的话，最好是md，如果不是的话，txt的不会乱")
    parser.add_argument("--language", type=str, default='zh', help="The other output lauguage is English, is en")
    import time

    # loop file in folder
    args = parser.parse_args()
    if args.sort == 'Relevance':
        sort = arxiv.SortCriterion.Relevance
    elif args.sort == 'LastUpdatedDate':
        sort = arxiv.SortCriterion.LastUpdatedDate
    else:
        sort = arxiv.SortCriterion.Relevance
    reader1 = Reader(key_word=args.key_word,
                     query=args.query,
                     filter_keys=args.filter_keys,
                     sort=sort,
                     args=args
                     )
    reader1.show_info()

    paper_list = []
    if args.pdf_path.endswith(".pdf"):
        paper_list.append(Paper(path=args.pdf_path))
    else:
        for root, dirs, files in os.walk(args.pdf_path):
            print("root:", root, "dirs:", dirs, 'files:', files)  # 当前目录路径
            count = 0
            for filename in files:
                # 如果找到PDF文件，则将其复制到目标文件夹中
                if filename.endswith(".pdf"):
                    print("---processing",count,filename,"---")

                    cur_paper = Paper(path=os.path.join(root, filename))

                    start_time = time.time()
                    result,title = reader1.summary_with_chat(paper_list=[cur_paper])
                    # text = '\n'.join(result[0])
                    # result_dict, match_title = analyze_text(text)
                    reader1.export_to_markdown("\n".join(result[0]),
                                               file_name="export/" + title[0]+".md")
                    print("md filename", "export/" + title[0])
                    print("---FINISH summary:", count, title[0], time.time() - start_time)

                    start_time = time.time()
                    # paperTitle = reader1.validateTitle(cur_paper.title[:80])

                    # match_title = re.sub('[{}]'.format(zhon.hanzi.punctuation), "_", match_title)
                    images = extract_imgs_from_pdf(title[0], cur_paper.path)
                    reader1.export_to_markdown("\n".join(result[0]),
                                               file_name="export/" + title[0]+"/summary.md")
                    # print("result_dict",result_dict)
                    # np.save("export/result_dict.npy", result_dict)
                    # with open("export/result_"+paperTitle+".json", "w") as outfile:
                    #     json.dump(result_dict, outfile)
                    # print("images", images)
                    # np.save("export/images.npy", images)
                    #
                    # write_image_pptx("CHI2023_chatpaper", result_dict, images)
                    try:
                        copyfile(cur_paper.path, "export/" + title[0] + ".pdf")
                        move(cur_paper.path, "chi2023/finish/" + filename)
                    except IOError as e:
                        print("Unable to copy file. %s" % e)
                        # exit(1)
                    except:
                        print("Unexpected error:", sys.exc_info())
                        # exit(1)

                    paper_list.append(cur_paper)
                    print("title+images time:", time.time() - start_time)
                    # print("---finish processing", filename, cur_paper.title, "---\n\n")
                    #

                    count = count+1



    # for paper_index, paper in enumerate(paper_list):
    #     start_time = time.time()
    #
    #     result = reader1.summary_with_chat(paper_list=[paper])
    #     text = '\n'.join(result)
    #     result_dict, match_title = analyze_text(text)
    #     reader1.export_to_markdown("\n".join(result[0]), file_name="export/"+reader1.validateTitle(paper.title[:80]))
    #     print("md filename", "export/"+paperTitle)
    #
    #     print("summary time:", match_title, time.time() - start_time)
